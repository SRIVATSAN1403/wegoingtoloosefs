import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)

    # Design System Colors (Steep Editorial Light Theme)
    COLOR_INK_BLACK = RGBColor(23, 25, 28)
    COLOR_PAPER_WHITE = RGBColor(255, 255, 255)
    COLOR_MIST_GRAY = RGBColor(242, 242, 243)
    COLOR_SLATE_GRAY = RGBColor(119, 123, 134)
    COLOR_BLUSH_PEACH = RGBColor(251, 225, 209)
    COLOR_SIENNA_BROWN = RGBColor(93, 42, 26)

    # Helper function to style text frame
    def style_title(shape, text, size=40, font_name="Georgia", color=COLOR_INK_BLACK):
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(size)
        p.font.bold = False
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.LEFT

    def add_bullet(tf, text, size=18, font_name="Arial", color=COLOR_INK_BLACK, bold=False):
        p = tf.add_paragraph()
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(10)

    # --- SLIDE 1: Title Slide (Accent Blush Peach background) ---
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape for Blush Peach theme
    bg = slide.shapes.add_shape(
        1, # Rectangle
        0, 0, Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BLUSH_PEACH
    bg.line.fill.background()

    # Title & Subtitle box
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "PyroTrace"
    p.font.name = "Georgia"
    p.font.size = Pt(64)
    p.font.color.rgb = COLOR_SIENNA_BROWN
    p.space_after = Pt(12)

    p2 = tf.add_paragraph()
    p2.text = "\"The AI that finds the frayed wire, not just the fire.\""
    p2.font.name = "Georgia"
    p2.font.size = Pt(28)
    p2.font.italic = True
    p2.font.color.rgb = COLOR_SIENNA_BROWN
    p2.space_after = Pt(24)

    p3 = tf.add_paragraph()
    p3.text = "Build Beyond Boundaries"
    p3.font.name = "Arial"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_SIENNA_BROWN

    # --- SLIDE 2: The Problem (Paper White background) ---
    slide = prs.slides.add_slide(slide_layout)
    # Background
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PAPER_WHITE
    bg.line.fill.background()

    # Title
    t_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.33), Inches(1.2))
    style_title(t_box, "The Problem: Alert Storms & Hidden Root Causes", size=36)

    # Content Columns
    left_col = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.5))
    tf_l = left_col.text_frame
    tf_l.word_wrap = True
    tf_l.paragraphs[0].text = "Standard Observability Gaps:"
    tf_l.paragraphs[0].font.bold = True
    tf_l.paragraphs[0].font.size = Pt(20)
    tf_l.paragraphs[0].font.name = "Georgia"
    add_bullet(tf_l, "• AI as a simple smoke detector: Flags anomaly thresholds but fails to trace the initial failure trigger.")
    add_bullet(tf_l, "• High MTTR (Mean Time to Resolution): Teams spend 80% of incident response time finding the root cause and only 20% fixing it.")

    right_col = slide.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.3), Inches(4.5))
    tf_r = right_col.text_frame
    tf_r.word_wrap = True
    tf_r.paragraphs[0].text = "Edge-Level Gaps:"
    tf_r.paragraphs[0].font.bold = True
    tf_r.paragraphs[0].font.size = Pt(20)
    tf_r.paragraphs[0].font.name = "Georgia"
    add_bullet(tf_r, "• IoT & Factory floor sensors stream critical telemetry to cloud servers, introducing latency and security compliance risks.")
    add_bullet(tf_r, "• Heavy Bayesian diagnostics models exceed edge compute capabilities.")

    # --- SLIDE 3: The Solution (Paper White background) ---
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PAPER_WHITE
    bg.line.fill.background()

    t_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.33), Inches(1.2))
    style_title(t_box, "The Solution: Edge-Native Causal Observation", size=36)

    left_col = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.5))
    tf_l = left_col.text_frame
    tf_l.word_wrap = True
    tf_l.paragraphs[0].text = "Edge Telemetry Pipeline:"
    tf_l.paragraphs[0].font.bold = True
    tf_l.paragraphs[0].font.size = Pt(20)
    tf_l.paragraphs[0].font.name = "Georgia"
    add_bullet(tf_l, "• ESP32-C3 Super Mini Streams real-time physical parameters (Fan RPM, Temperature, CPU Load) offline over local USB COM15.")
    add_bullet(tf_l, "• Rolling 60-second Pandas memory buffer handles data frequency spikes with flat, stable memory utilization.")

    right_col = slide.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.3), Inches(4.5))
    tf_r = right_col.text_frame
    tf_r.word_wrap = True
    tf_r.paragraphs[0].text = "AI Causal Engine:"
    tf_r.paragraphs[0].font.bold = True
    tf_r.paragraphs[0].font.size = Pt(20)
    tf_r.paragraphs[0].font.name = "Georgia"
    add_bullet(tf_r, "• Real-Time Outlier Tagging: Scikit-Learn Isolation Forest highlights statistical telemetry breaches instantly.")
    add_bullet(tf_r, "• Causal Sequence Detector: Traces time-lagged cross-correlation backward through data to identify initial physical trigger.")

    # --- SLIDE 4: Competitive Landscape (Paper White background) ---
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PAPER_WHITE
    bg.line.fill.background()

    t_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.33), Inches(1.2))
    style_title(t_box, "Competitive Observability Landscape", size=36)

    content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "Positioning vs. Enterprise Competitors:"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.name = "Georgia"

    add_bullet(tf, "• Dynatrace (Davis AI): Cloud-scale app mapping costing millions. PyroTrace targets localized, low-cost Edge & IoT hardware.")
    add_bullet(tf, "• Datadog (Watchdog AI): Cloud-dependent SaaS telemetry. PyroTrace operates 100% offline, keeping sensitive factory data on-premises.")
    add_bullet(tf, "• BigPanda: Software ticketing event correlation. PyroTrace correlates physical physics (RPM, temperature, motor load).")

    # --- SLIDE 5: Value Proposition & Impact (Blush Peach background) ---
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BLUSH_PEACH
    bg.line.fill.background()

    t_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.33), Inches(1.2))
    style_title(t_box, "Value Proposition & Business Impact", size=36, color=COLOR_SIENNA_BROWN)

    left_col = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.5))
    tf_l = left_col.text_frame
    tf_l.word_wrap = True
    tf_l.paragraphs[0].text = "Financial & Operational Gains:"
    tf_l.paragraphs[0].font.bold = True
    tf_l.paragraphs[0].font.size = Pt(20)
    tf_l.paragraphs[0].font.name = "Georgia"
    tf_l.paragraphs[0].font.color.rgb = COLOR_SIENNA_BROWN
    add_bullet(tf_l, "• 80% MTTR Reduction: Eliminates investigation search time by mapping root cause triggers instantly.", color=COLOR_SIENNA_BROWN)
    add_bullet(tf_l, "• Cost Efficiency: Catch hardware failures (e.g. fan bearings) seconds before thermal shutdown saves costly components.", color=COLOR_SIENNA_BROWN)

    right_col = slide.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.3), Inches(4.5))
    tf_r = right_col.text_frame
    tf_r.word_wrap = True
    tf_r.paragraphs[0].text = "Strategic Advantages:"
    tf_r.paragraphs[0].font.bold = True
    tf_r.paragraphs[0].font.size = Pt(20)
    tf_r.paragraphs[0].font.name = "Georgia"
    tf_r.paragraphs[0].font.color.rgb = COLOR_SIENNA_BROWN
    add_bullet(tf_r, "• 100% Data Sovereignty: Zero cloud dependency meets government, defense, and power plant security standards.", color=COLOR_SIENNA_BROWN)
    add_bullet(tf_r, "• Team Empowerment: Democratizes debugging, giving junior technicians senior diagnostic capabilities.", color=COLOR_SIENNA_BROWN)

    # Save to docs folder
    os.makedirs("docs", exist_ok=True)
    save_path = os.path.join("docs", "NEURAL_X_PPT.pptx")
    prs.save(save_path)
    print(f"Presentation slide deck created successfully at: {save_path}")

if __name__ == "__main__":
    create_presentation()
